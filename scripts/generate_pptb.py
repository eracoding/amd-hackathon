"""Generate PPT-B: the original presentation (PPT-A) enriched with what the
room did to it — detected on-slide annotations re-drawn in place, a discreet
per-slide capture panel (engagement, questions, annotations), and one final
synthesis slide: "what matters for you, as a team".

    python -m scripts.generate_pptb recordings/dryrun1/events.jsonl \
        --pptx deck.pptx               # enrich the original PowerPoint
    python -m scripts.generate_pptb sessions/session.jsonl \
        --deck slides.pdf              # no .pptx? rebuild slides from the PDF

Inputs: the session/ingested event stream (JSONL). Annotation patches
referenced by ScreenAnnotationEvent.patch_path are pasted at their detected
position. If reports/last_debrief.json exists, its summary/action items feed
the synthesis slide; otherwise a deterministic synthesis is computed from
the events alone.
"""
from __future__ import annotations

import argparse
import json
import sys
from collections import defaultdict
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation                      # noqa: E402
from pptx.dml.color import RGBColor                # noqa: E402
from pptx.util import Emu, Pt                      # noqa: E402

INK = RGBColor(0xE2, 0x57, 0x4C)        # panel accent
MUTED = RGBColor(0x6B, 0x74, 0x80)
DARK = RGBColor(0x20, 0x26, 0x30)


def load_events(path: Path) -> list[dict]:
    return [json.loads(ln) for ln in
            path.read_text(encoding="utf-8").splitlines() if ln.strip()]


def aggregate(events: list[dict]) -> dict:
    """Per-slide engagement / questions / annotations + global rollups."""
    slides: dict[int, dict] = defaultdict(
        lambda: {"att_sum": 0.0, "att_n": 0, "questions": [],
                 "annotations": [], "title": ""})
    current = 0
    unresolved: list[str] = []
    for e in events:
        t = e["topic"]
        if t == "SlideChange":
            current = e["slide"]
            slides[current]["title"] = e.get("title", "")
        elif t == "AttentionEvent":
            s = slides[current]
            s["att_sum"] += e["attention"]
            s["att_n"] += 1
        elif t == "InteractionEvent" and e.get("kind") == "question":
            n = e.get("slide") or current
            slides[n]["questions"].append(
                {"who": e.get("person_id", "?"), "text": e.get("text", ""),
                 "src": e.get("source", "typed")})
            if not e.get("resolved", False):
                unresolved.append(e.get("text", ""))
        elif t == "ScreenAnnotationEvent":
            slides[e["slide"]]["annotations"].append(e)
    for n, s in slides.items():
        s["engagement"] = (s["att_sum"] / s["att_n"]) if s["att_n"] else None
    return {"slides": dict(slides), "unresolved": unresolved}


def _panel(slide_obj, sw: int, sh: int, lines: list[str],
           ann_bboxes: list[list[float]] | None = None) -> None:
    """Discreet capture panel; picks a corner free of annotations."""
    w = int(sw * 0.30)
    h = Pt(13 + 13 * len(lines))
    corners = [(sw - w, sh - int(h)), (0, sh - int(h)),   # BR, BL
               (sw - w, 0), (0, 0)]                        # TR, TL
    def overlaps(cx, cy):
        px0, py0 = cx / sw, cy / sh
        px1, py1 = (cx + w) / sw, (cy + int(h)) / sh
        return any(b[0] < px1 and px0 < b[2] and b[1] < py1 and py0 < b[3]
                   for b in (ann_bboxes or []))
    cx, cy = next(((x, y) for x, y in corners if not overlaps(x, y)),
                  corners[0])
    pad = Emu(91440 // 4)
    box = slide_obj.shapes.add_textbox(
        max(pad, min(cx, sw - w - pad)),
        max(pad, min(cy, sh - int(h) - pad)), w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(6)
    tf.margin_top = tf.margin_bottom = Pt(4)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.font.size = Pt(9)
        p.font.color.rgb = INK if i == 0 else MUTED
        p.font.bold = i == 0


def _paste_annotations(slide_obj, sw: int, sh: int, anns: list[dict]) -> None:
    for a in anns:
        p = a.get("patch_path", "")
        x0, y0, x1, y1 = a["bbox"]
        left, top = int(x0 * sw), int(y0 * sh)
        w, h = max(1, int((x1 - x0) * sw)), max(1, int((y1 - y0) * sh))
        if p and Path(p).exists():
            slide_obj.shapes.add_picture(p, left, top, width=w, height=h)
        else:  # no patch (simulated/legacy events): mark the region
            box = slide_obj.shapes.add_textbox(left, top, w, max(h, Pt(14)))
            tf = box.text_frame
            tf.text = "✎ annotation"
            tf.paragraphs[0].font.size = Pt(10)
            tf.paragraphs[0].font.color.rgb = INK


def _synthesis_slide(prs: Presentation, agg: dict, debrief: dict) -> None:
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 \
        else prs.slide_layouts[-1]                    # blank-ish layout
    s = prs.slides.add_slide(layout)
    sw, sh = prs.slide_width, prs.slide_height
    title = s.shapes.add_textbox(int(sw * .06), int(sh * .07),
                                 int(sw * .88), int(sh * .14))
    p = title.text_frame.paragraphs[0]
    p.text = "What matters for you, as a team"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK

    body = s.shapes.add_textbox(int(sw * .06), int(sh * .24),
                                int(sw * .88), int(sh * .68))
    tf = body.text_frame
    tf.word_wrap = True
    lines: list[tuple[str, bool]] = []
    if debrief.get("summary"):
        lines.append((debrief["summary"], False))
    eng = [(n, v["engagement"]) for n, v in agg["slides"].items()
           if v["engagement"] is not None and n > 0]
    if eng:
        lo = min(eng, key=lambda t: t[1])
        ann_total = sum(len(v["annotations"]) for v in agg["slides"].values())
        q_total = sum(len(v["questions"]) for v in agg["slides"].values())
        if lo[1] < 0.75:
            lines.append((f"Attention converged least on slide {lo[0]} "
                          f"({lo[1]:.0%} engagement) — revisit it.", False))
        else:
            lines.append((f"Attention stayed high throughout "
                          f"(lowest: slide {lo[0]} at {lo[1]:.0%}).", False))
        lines.append((f"The team produced {ann_total} on-slide annotation(s) "
                      f"and {q_total} question(s).", False))
    if agg["unresolved"]:
        lines.append(("Unresolved questions:", True))
        lines += [(f"• {q}", False) for q in agg["unresolved"][:5]]
    for item in debrief.get("action_items", [])[:4]:
        lines.append((f"→ {item}", False))
    for i, (txt, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(16)
        p.font.bold = bold
        p.font.color.rgb = DARK


def build_from_pdf(deck_pdf: Path, n_slides: int) -> Presentation:
    """No original .pptx: rebuild a deck with the PDF pages as full-bleed
    slide images, so annotations land on faithful backgrounds."""
    from aura.perception.slides import DeckIndex
    deck = DeckIndex(deck_pdf)
    prs = Presentation()
    img0 = deck.render_page(1, scale=2.0)
    prs.slide_width = Emu(9144000)                       # 10in 16:9-ish
    prs.slide_height = int(prs.slide_width * img0.height / img0.width)
    blank = prs.slide_layouts[6]
    import io
    for n in range(1, min(n_slides, len(deck)) + 1):
        s = prs.slides.add_slide(blank)
        buf = io.BytesIO()
        deck.render_page(n, scale=2.0).save(buf, format="PNG")
        buf.seek(0)
        s.shapes.add_picture(buf, 0, 0, width=prs.slide_width,
                             height=prs.slide_height)
    return prs


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("events", help="session/ingested events JSONL")
    ap.add_argument("--pptx", help="original PPT-A (.pptx) to enrich")
    ap.add_argument("--deck", help="PDF fallback when no .pptx is available")
    ap.add_argument("--debrief", default="",
                help="optional debrief JSON from THIS session")
    ap.add_argument("--out", default="reports/PPT_B.pptx")
    args = ap.parse_args()

    events = load_events(Path(args.events))
    agg = aggregate(events)
    debrief = {}
    if args.debrief and Path(args.debrief).exists():
        debrief = json.loads(Path(args.debrief).read_text())

    max_slide = max((n for n in agg["slides"] if n > 0), default=0)
    if args.pptx:
        prs = Presentation(args.pptx)
    elif args.deck:
        prs = build_from_pdf(Path(args.deck), max_slide)
    else:
        sys.exit("provide --pptx (original deck) or --deck (PDF fallback)")

    sw, sh = prs.slide_width, prs.slide_height
    for i, slide_obj in enumerate(prs.slides, start=1):
        s = agg["slides"].get(i)
        if not s:
            continue
        _paste_annotations(slide_obj, sw, sh, s["annotations"])
        lines = ["AURA capture"]
        if s["engagement"] is not None:
            lines.append(f"engagement {s['engagement']:.0%}")
        if s["questions"]:
            lines.append(f"{len(s['questions'])} question(s): "
                         + "; ".join(q['text'][:40] for q in s['questions'][:2]))
        if s["annotations"]:
            lines.append(f"{len(s['annotations'])} annotation(s) drawn here")
        if len(lines) > 1:
            _panel(slide_obj, sw, sh, lines,
                   ann_bboxes=[a["bbox"] for a in s["annotations"]])

    _synthesis_slide(prs, agg, debrief)
    out = Path(args.out)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"PPT-B written -> {out} "
          f"({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
