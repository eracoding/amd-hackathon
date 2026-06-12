"""PPT-B generator test: events + original deck -> enriched .pptx."""
import json
import subprocess
import sys
from pathlib import Path

import pytest

pytest.importorskip("pptx")
from pptx import Presentation               # noqa: E402
from pptx.util import Emu, Pt               # noqa: E402

from aura.events import (                   # noqa: E402
    AttentionEvent, InteractionEvent, InteractionKind,
    ScreenAnnotationEvent, SessionEnd, SlideChange,
)


def _write_events(path: Path) -> None:
    evs = []
    t = 1000.0
    for slide in (1, 2, 3):
        e = SlideChange(slide=slide, title=f"Slide {slide}"); e.ts = t
        evs.append(e)
        for k in range(4):
            a = AttentionEvent(person_id="person_0",
                               attention=0.4 if slide == 2 else 0.9)
            a.ts = t + 0.2 * k
            evs.append(a)
        t += 2
    q = InteractionEvent(person_id="voice", kind=InteractionKind.question,
                         text="Why slide two?", slide=2, source="voice")
    q.ts = 1003; evs.append(q)
    ann = ScreenAnnotationEvent(slide=2, bbox=[0.4, 0.4, 0.55, 0.5],
                                area_frac=0.004, kind="drawn")
    ann.ts = 1003.5; evs.append(ann)
    end = SessionEnd(); end.ts = t; evs.append(end)
    with path.open("w") as fh:
        for e in evs:
            fh.write(json.dumps({"topic": e.topic, **e.model_dump()}) + "\n")


def _write_ppta(path: Path) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    for n in (1, 2, 3):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = f"Slide {n}"
    prs.save(path)


def test_pptb_generation(tmp_path):
    events = tmp_path / "events.jsonl"
    ppta = tmp_path / "A.pptx"
    out = tmp_path / "B.pptx"
    _write_events(events)
    _write_ppta(ppta)
    r = subprocess.run(
        [sys.executable, "-m", "scripts.generate_pptb", str(events),
         "--pptx", str(ppta), "--out", str(out)],
        capture_output=True, text=True, cwd=Path(__file__).parent.parent)
    assert r.returncode == 0, r.stderr
    prs = Presentation(out)
    assert len(prs.slides) == 4, "3 originals + 1 synthesis"
    texts = ["\n".join(sh.text_frame.text for sh in sl.shapes
                       if sh.has_text_frame) for sl in prs.slides]
    assert "AURA capture" in texts[1], "slide 2 must carry the capture panel"
    assert "annotation" in texts[1]
    assert "What matters for you, as a team" in texts[3]
    assert "Why slide two?" in texts[3], "unresolved question in synthesis"
    assert "slide 2" in texts[3], "dip slide named in synthesis"
